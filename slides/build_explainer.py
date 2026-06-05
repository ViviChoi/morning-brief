"""Build the Morning Brief explainer deck.

Run:
    .venv/bin/python slides/build_explainer.py

Output:
    slides/MorningBrief-Explainer.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
C_BG = RGBColor(0x2D, 0x37, 0x48)        # Title area: dark blue-grey
C_BG_LIGHT = RGBColor(0xF7, 0xF8, 0xFA)  # Slide body: off-white
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN = RGBColor(0x1A, 0x7F, 0x37)     # Accent: growth / positive
C_RED = RGBColor(0xCF, 0x22, 0x2E)       # Accent: loss / danger
C_AMBER = RGBColor(0xD4, 0x8A, 0x00)     # Warning
C_BLUE_MID = RGBColor(0x3B, 0x82, 0xF6)  # Sub-accent
C_GRAY_TEXT = RGBColor(0x44, 0x44, 0x55) # Body text
C_GRAY_DIM = RGBColor(0x99, 0x99, 0xAA)  # Dim text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).parent / "MorningBrief-Explainer.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line_color: RGBColor | None = None, line_width_pt: float = 0.0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text: str, left, top, width, height,
                 font_size: int = 20,
                 bold: bool = False,
                 color: RGBColor = C_GRAY_TEXT,
                 align=PP_ALIGN.LEFT,
                 wrap: bool = True,
                 font_name: str = "Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines: list[str | tuple], left, top, width, height,
                  default_size: int = 18,
                  default_color: RGBColor = C_GRAY_TEXT,
                  font_name: str = "Calibri",
                  line_spacing_pt: float | None = None):
    """Add a text box with multiple paragraphs.

    Each item in `lines` is either:
      - a plain str  → rendered with defaults
      - a tuple (text, size, bold, color)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color = item, default_size, False, default_color
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
        if line_spacing_pt:
            from pptx.util import Pt as _Pt
            p.line_spacing = _Pt(line_spacing_pt)
    return txBox


def header_band(slide, title: str, subtitle: str = ""):
    """Dark header band across the top ~30% of the slide."""
    band_h = Inches(2.2)
    add_rect(slide, 0, 0, SLIDE_W, band_h, fill=C_BG)
    # Title
    add_text_box(slide, title,
                 left=Inches(0.6), top=Inches(0.35),
                 width=Inches(12.0), height=Inches(1.1),
                 font_size=38, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     left=Inches(0.6), top=Inches(1.45),
                     width=Inches(12.0), height=Inches(0.55),
                     font_size=18, bold=False, color=C_GRAY_DIM,
                     align=PP_ALIGN.LEFT)


def slide_number(slide, n: int):
    add_text_box(slide, str(n),
                 left=Inches(12.6), top=Inches(7.1),
                 width=Inches(0.6), height=Inches(0.3),
                 font_size=11, color=C_GRAY_DIM,
                 align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide factories
# ---------------------------------------------------------------------------

def slide_01_title(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG)

    # Full background gradient feel — just a subtle lighter rect at bottom
    add_rect(slide, 0, Inches(5.0), SLIDE_W, Inches(2.5), fill=RGBColor(0x1A, 0x20, 0x2E))

    # Green accent bar on left edge
    add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, fill=C_GREEN)

    add_text_box(slide,
                 "Morning Brief",
                 left=Inches(0.55), top=Inches(1.1),
                 width=Inches(12.0), height=Inches(1.0),
                 font_size=52, bold=True, color=C_WHITE)

    add_text_box(slide,
                 "一个跑在你桌面上的市场情报系统",
                 left=Inches(0.55), top=Inches(2.15),
                 width=Inches(12.0), height=Inches(0.8),
                 font_size=30, bold=False, color=RGBColor(0xA0, 0xC4, 0xFF))

    add_text_box(slide,
                 "为不会读金融代码的人设计  ·  跨 Mac + Win + iPhone  ·  自包含、可分享",
                 left=Inches(0.55), top=Inches(3.0),
                 width=Inches(12.0), height=Inches(0.6),
                 font_size=18, bold=False, color=C_GRAY_DIM)

    # Tag chips
    for i, (tag, color) in enumerate([
        ("全免费数据源", C_GREEN),
        ("AI 白话解读", C_BLUE_MID),
        ("本地隐私保护", C_AMBER),
    ]):
        x = Inches(0.55 + i * 3.1)
        add_rect(slide, x, Inches(4.0), Inches(2.8), Inches(0.5),
                 fill=color, line_color=None)
        add_text_box(slide, tag,
                     left=x + Inches(0.05), top=Inches(4.03),
                     width=Inches(2.7), height=Inches(0.45),
                     font_size=16, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "今天是 2026-06-05  |  Phase 4 shipped  |  35 tests passing",
                 left=Inches(0.55), top=Inches(6.85),
                 width=Inches(12.0), height=Inches(0.4),
                 font_size=12, color=C_GRAY_DIM)


def slide_02_problem(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "你早上想知道市场咋样，但是…", "问题陈述")

    pain_points = [
        ("📱  主流财经 APP 推送太多太乱", "彭博、雪球、东财…推送 100 条，你真正需要的 5 条被淹没了"),
        ("💸  Day1Global / 彭博终端", "专业工具看不懂，或者太贵（Bloomberg 终端 ~$2,000/月）"),
        ("🗂️  持仓散在多个 App", "IBKR + 币安 + Coinbase，没有任何一个工具帮你串起来看全局"),
    ]

    for i, (title, desc) in enumerate(pain_points):
        y = Inches(2.5 + i * 1.45)
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.25),
                 fill=C_WHITE, line_color=RGBColor(0xDD, 0xDD, 0xEE), line_width_pt=1.0)
        add_text_box(slide, title,
                     left=Inches(0.75), top=y + Inches(0.1),
                     width=Inches(11.8), height=Inches(0.5),
                     font_size=20, bold=True, color=C_BG)
        add_text_box(slide, desc,
                     left=Inches(0.75), top=y + Inches(0.58),
                     width=Inches(11.8), height=Inches(0.5),
                     font_size=16, color=C_GRAY_TEXT)

    slide_number(slide, 2)


def slide_03_solution(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "解决方案：一个文件夹，两个工具", "设计原则")

    # Two surface cards
    for i, (color, label, tagline, bullets) in enumerate([
        (C_GREEN, "Surface A — Daily Brief", "每天跑一次  ·  5 分钟知道市场", [
            "抓取 9 个核心市场指标",
            "BTC 抄底分（0–100）",
            "Claude 白话 AI 解读",
            "生成可分享 HTML 快照",
        ]),
        (C_BLUE_MID, "Surface B — Portfolio Explorer", "本地 WebUI  ·  上传 CSV 看持仓", [
            "IBKR / Binance / Coinbase 导入",
            "持仓饼图 + 盈亏棒图",
            "AI 白话总结（4 个板块）",
            "只跑本地，CSV 不外泄",
        ]),
    ]):
        x = Inches(0.4 + i * 6.45)
        add_rect(slide, x, Inches(2.3), Inches(6.1), Inches(4.8),
                 fill=C_WHITE, line_color=color, line_width_pt=2.5)
        add_rect(slide, x, Inches(2.3), Inches(6.1), Inches(0.7), fill=color)
        add_text_box(slide, label,
                     left=x + Inches(0.15), top=Inches(2.32),
                     width=Inches(5.8), height=Inches(0.6),
                     font_size=18, bold=True, color=C_WHITE)
        add_text_box(slide, tagline,
                     left=x + Inches(0.15), top=Inches(3.1),
                     width=Inches(5.8), height=Inches(0.4),
                     font_size=14, color=C_GRAY_DIM)
        for j, b in enumerate(bullets):
            add_text_box(slide, "▸  " + b,
                         left=x + Inches(0.2), top=Inches(3.6 + j * 0.72),
                         width=Inches(5.7), height=Inches(0.6),
                         font_size=17, color=C_GRAY_TEXT)

    slide_number(slide, 3)


def slide_04_surface_a(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "Surface A — Daily Brief", "每天早上 8:00 自动跑一次（Mac mini launchd）")

    items = [
        ("9 个核心指标", "SPY / QQQ / VIX / DXY / GLD / USO / BTC / ETH — 大盘、情绪、宏观全覆盖"),
        ("Crypto F&G 指数", "alternative.me 每日发布  ·  0 = 极恐  ·  100 = 极贪"),
        ("BTC 抄底分  0–100", "4 个免费指标加权：ETF 资金流 + Funding Rate + F&G + 多空比"),
        ("Claude AI 5 段解读", "Headline · 流动性基调 · 各票 thesis · BTC 叙事 · Top 10 要闻"),
        ("自包含 HTML 快照", "demos/YYYY-MM-DD.html  — 发给任何人，浏览器直接打开，零安装"),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(2.4 + i * 0.95)
        # Number badge
        add_rect(slide, Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.45),
                 fill=C_GREEN)
        add_text_box(slide, str(i + 1),
                     left=Inches(0.4), top=y + Inches(0.06),
                     width=Inches(0.45), height=Inches(0.4),
                     font_size=16, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, title,
                     left=Inches(1.1), top=y,
                     width=Inches(3.5), height=Inches(0.5),
                     font_size=18, bold=True, color=C_BG)
        add_text_box(slide, desc,
                     left=Inches(1.1), top=y + Inches(0.45),
                     width=Inches(11.2), height=Inches(0.4),
                     font_size=15, color=C_GRAY_TEXT)

    # Cost note
    add_rect(slide, Inches(9.5), Inches(2.4), Inches(3.4), Inches(1.1),
             fill=RGBColor(0xE6, 0xF4, 0xEA), line_color=C_GREEN, line_width_pt=1.5)
    add_text_box(slide, "💰  成本",
                 left=Inches(9.65), top=Inches(2.45),
                 width=Inches(3.1), height=Inches(0.4),
                 font_size=15, bold=True, color=C_GREEN)
    add_text_box(slide, "数据全免费  +  Claude API ~$0.05/天",
                 left=Inches(9.65), top=Inches(2.82),
                 width=Inches(3.1), height=Inches(0.5),
                 font_size=14, color=C_GRAY_TEXT)

    slide_number(slide, 4)


def slide_05_surface_b(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "Surface B — Portfolio Explorer", "本地 WebUI  ·  http://127.0.0.1:8765  ·  只有你自己能访问")

    # Left: steps
    steps = [
        ("① 上传 CSV", "IBKR / Binance / Coinbase 原生格式\n或手工表（symbol + quantity + cost 三列）"),
        ("② 看可视化", "持仓饼图（Allocation Donut）+ 盈亏棒图（P&L Bar）"),
        ("③ 读 AI 解读", "Claude 写 4 段白话：你是什么风格、什么在赚、什么在亏、要警惕啥"),
    ]
    for i, (title, desc) in enumerate(steps):
        y = Inches(2.5 + i * 1.5)
        add_rect(slide, Inches(0.4), y, Inches(5.8), Inches(1.3),
                 fill=C_WHITE, line_color=C_BLUE_MID, line_width_pt=2.0)
        add_text_box(slide, title,
                     left=Inches(0.6), top=y + Inches(0.1),
                     width=Inches(5.4), height=Inches(0.5),
                     font_size=19, bold=True, color=C_BLUE_MID)
        add_text_box(slide, desc,
                     left=Inches(0.6), top=y + Inches(0.55),
                     width=Inches(5.4), height=Inches(0.65),
                     font_size=14, color=C_GRAY_TEXT)

    # Right: privacy box
    add_rect(slide, Inches(7.0), Inches(2.5), Inches(5.9), Inches(4.5),
             fill=RGBColor(0xFF, 0xF8, 0xE1), line_color=C_AMBER, line_width_pt=2.0)
    add_text_box(slide, "🔒  隐私保证",
                 left=Inches(7.2), top=Inches(2.6),
                 width=Inches(5.5), height=Inches(0.55),
                 font_size=20, bold=True, color=C_AMBER)
    privacy_points = [
        "服务器只绑定 127.0.0.1（本机回环地址）",
        "外网完全无法访问，即使在公共 WiFi",
        "上传的 CSV 存在 portfolios/ 文件夹",
        "portfolios/ 被 .gitignore 排除",
        "即使你 git push，CSV 绝对不会上传",
        "Anthropic API 不允许用调用数据训练模型",
    ]
    for i, pt in enumerate(privacy_points):
        add_text_box(slide, "✓  " + pt,
                     left=Inches(7.2), top=Inches(3.25 + i * 0.58),
                     width=Inches(5.5), height=Inches(0.5),
                     font_size=15, color=C_GRAY_TEXT)

    slide_number(slide, 5)


def slide_06_btc_score(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "BTC 抄底分  0–100 — 怎么打出来的", "4 个免费日线指标加权合成  ·  借鉴 Day1Global 方法论")

    # Four component cards
    components = [
        (C_GREEN, "ETF 净资金流", "38%", "farside.co.uk", "机构每天往 BTC 现货 ETF\n净买入还是净卖出？\n净流入 +500M USD → 满分 100"),
        (C_RED, "Funding Rate", "25%", "Binance Futures API", "永续合约多头付给空头的费率\n费率极正 = 多头过热\n费率极负 = 空头恐慌"),
        (C_AMBER, "F&G 指数", "22%", "alternative.me", "0–100 的综合情绪指标\n直接映射为组件分\n0 = 极恐  100 = 极贪"),
        (C_BLUE_MID, "多空比", "15%", "Binance Futures API", "合约市场做多账户 ÷ 做空账户\n> 1 = 偏多\n< 1 = 偏空"),
    ]

    for i, (color, title, weight, source, desc) in enumerate(components):
        x = Inches(0.3 + i * 3.25)
        add_rect(slide, x, Inches(2.35), Inches(3.1), Inches(3.8),
                 fill=C_WHITE, line_color=color, line_width_pt=2.5)
        add_rect(slide, x, Inches(2.35), Inches(3.1), Inches(0.65), fill=color)
        add_text_box(slide, title,
                     left=x + Inches(0.1), top=Inches(2.38),
                     width=Inches(2.9), height=Inches(0.55),
                     font_size=16, bold=True, color=C_WHITE)
        add_text_box(slide, "权重 " + weight,
                     left=x + Inches(0.1), top=Inches(3.1),
                     width=Inches(2.9), height=Inches(0.4),
                     font_size=22, bold=True, color=color)
        add_text_box(slide, source,
                     left=x + Inches(0.1), top=Inches(3.55),
                     width=Inches(2.9), height=Inches(0.35),
                     font_size=12, color=C_GRAY_DIM)
        add_text_box(slide, desc,
                     left=x + Inches(0.1), top=Inches(3.95),
                     width=Inches(2.9), height=Inches(1.1),
                     font_size=13, color=C_GRAY_TEXT)

    # Scale bar at bottom
    scale_labels = [
        (0, "0", "极度恐惧\n历史买点", C_RED),
        (1, "15", "", C_AMBER),
        (2, "45", "恐惧", C_AMBER),
        (3, "55", "中性", C_GREEN),
        (4, "85", "贪婪", C_AMBER),
        (5, "100", "极度贪婪\n历史卖点", C_RED),
    ]
    bar_y = Inches(6.3)
    bar_left = Inches(0.5)
    bar_w = Inches(12.3)
    bar_h = Inches(0.35)
    add_rect(slide, bar_left, bar_y, bar_w, bar_h,
             fill=RGBColor(0xDD, 0xDD, 0xEE))
    # Color segments
    segs = [
        (0.0, 0.15, C_RED),
        (0.15, 0.45, C_AMBER),
        (0.45, 0.55, C_GREEN),
        (0.55, 0.85, C_AMBER),
        (0.85, 1.0, C_RED),
    ]
    for s, e, c in segs:
        add_rect(slide,
                 bar_left + Emu(int(bar_w * s)),
                 bar_y,
                 Emu(int(bar_w * (e - s))),
                 bar_h,
                 fill=c)
    # Labels
    for frac, num, label, color in scale_labels:
        x_pos = bar_left + Emu(int(bar_w * (int(num) / 100)))
        add_text_box(slide, num,
                     left=x_pos - Inches(0.25), top=bar_y - Inches(0.35),
                     width=Inches(0.5), height=Inches(0.3),
                     font_size=11, bold=True, color=color,
                     align=PP_ALIGN.CENTER)
    slide_number(slide, 6)


def slide_07_data_sources(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "数据源 — 全部免费，无需付费订阅", "只需要 Anthropic API Key（~$0.05/天）")

    sources = [
        (C_GREEN, "yfinance", "美股 / ETF / VIX / DXY / Gold / Oil", "免费  ·  Yahoo Finance 公共 API"),
        (C_BLUE_MID, "CoinGecko", "加密货币现货价格 + 24h 涨跌幅", "免费  ·  无需注册"),
        (C_AMBER, "alternative.me", "Crypto Fear & Greed Index", "免费  ·  每日更新"),
        (C_RED, "farside.co.uk", "BTC 现货 ETF 每日净资金流数据", "免费  ·  公开网页抓取"),
        (C_BLUE_MID, "Binance Futures API", "Funding Rate + 多空比", "免费  ·  无需账户"),
        (C_GREEN, "Anthropic Claude", "所有 AI 白话解读文字", "~$0.05/天  ·  需要 API Key"),
    ]

    for i, (color, name, usage, note) in enumerate(sources):
        row = i % 3
        col = i // 3
        x = Inches(0.4 + col * 6.45)
        y = Inches(2.35 + row * 1.6)
        add_rect(slide, x, y, Inches(6.1), Inches(1.4),
                 fill=C_WHITE, line_color=color, line_width_pt=2.0)
        add_rect(slide, x, y, Inches(0.18), Inches(1.4), fill=color)
        add_text_box(slide, name,
                     left=x + Inches(0.28), top=y + Inches(0.1),
                     width=Inches(5.6), height=Inches(0.45),
                     font_size=18, bold=True, color=C_BG)
        add_text_box(slide, usage,
                     left=x + Inches(0.28), top=y + Inches(0.55),
                     width=Inches(5.6), height=Inches(0.4),
                     font_size=14, color=C_GRAY_TEXT)
        add_text_box(slide, note,
                     left=x + Inches(0.28), top=y + Inches(0.95),
                     width=Inches(5.6), height=Inches(0.35),
                     font_size=12, color=C_GRAY_DIM)

    slide_number(slide, 7)


def slide_08_architecture(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "技术架构 — 数据流一览", "单文件夹  ·  git 同步  ·  无需服务器  ·  可移植")

    # Node definitions: (label, x, y, w, h, fill_color)
    nodes = [
        ("config/\nwatchlist.json", Inches(0.3),  Inches(3.2),  Inches(2.2), Inches(0.9),  C_BLUE_MID),
        ("lib/fetch.py\n数据抓取",    Inches(3.0),  Inches(2.4),  Inches(2.3), Inches(0.9),  C_GREEN),
        ("lib/score.py\nBTC 抄底分",  Inches(3.0),  Inches(3.8),  Inches(2.3), Inches(0.9),  C_AMBER),
        ("lib/claude_brief.py\nAI 解读", Inches(6.1), Inches(2.4),  Inches(2.5), Inches(0.9),  C_RED),
        ("lib/render.py\nHTML 渲染",   Inches(6.1),  Inches(3.8),  Inches(2.5), Inches(0.9),  RGBColor(0x6F, 0x42, 0xC1)),
        ("snapshots/\nYYYY-MM-DD.json", Inches(9.3), Inches(2.4), Inches(2.5), Inches(0.9), C_BG),
        ("demos/\nYYYY-MM-DD.html",   Inches(9.3),  Inches(3.8),  Inches(2.5), Inches(0.9),  C_BG),
    ]

    for label, x, y, w, h, color in nodes:
        add_rect(slide, x, y, w, h, fill=color)
        add_text_box(slide, label,
                     left=x + Inches(0.08), top=y + Inches(0.05),
                     width=w - Inches(0.15), height=h - Inches(0.1),
                     font_size=13, bold=False, color=C_WHITE,
                     align=PP_ALIGN.CENTER)

    # Arrows (simple thin rects as lines)
    def h_arrow(x1, y1, x2):
        mid_y = y1 + Inches(0.45)
        add_rect(slide, x1, mid_y - Inches(0.02), x2 - x1, Inches(0.04), fill=C_GRAY_TEXT)

    def v_arrow(x1, y_top, y_bot):
        mid_x = x1 + Inches(0.02)
        add_rect(slide, mid_x, y_top, Inches(0.04), y_bot - y_top, fill=C_GRAY_TEXT)

    # watchlist → fetch
    h_arrow(Inches(2.5), Inches(2.85), Inches(3.0))
    # watchlist → score
    h_arrow(Inches(2.5), Inches(4.25), Inches(3.0))
    # fetch → claude
    h_arrow(Inches(5.3), Inches(2.85), Inches(6.1))
    # fetch → render
    v_arrow(Inches(4.1), Inches(3.3), Inches(3.8))
    h_arrow(Inches(5.3), Inches(4.25), Inches(6.1))
    # score → render
    h_arrow(Inches(5.3), Inches(4.25), Inches(6.1))
    # claude → snapshots
    h_arrow(Inches(8.6), Inches(2.85), Inches(9.3))
    # render → demos
    h_arrow(Inches(8.6), Inches(4.25), Inches(9.3))

    # Labels
    add_text_box(slide, "① 读 watchlist",
                 left=Inches(0.3), top=Inches(2.4), width=Inches(2.5), height=Inches(0.4),
                 font_size=12, color=C_GRAY_DIM)
    add_text_box(slide, "② 抓数据",
                 left=Inches(3.0), top=Inches(2.0), width=Inches(2.3), height=Inches(0.35),
                 font_size=12, color=C_GRAY_DIM)
    add_text_box(slide, "③ 算分数",
                 left=Inches(3.0), top=Inches(4.75), width=Inches(2.3), height=Inches(0.35),
                 font_size=12, color=C_GRAY_DIM)
    add_text_box(slide, "④ 生成文字",
                 left=Inches(6.1), top=Inches(2.0), width=Inches(2.5), height=Inches(0.35),
                 font_size=12, color=C_GRAY_DIM)
    add_text_box(slide, "⑤ 渲染 HTML",
                 left=Inches(6.1), top=Inches(4.75), width=Inches(2.5), height=Inches(0.35),
                 font_size=12, color=C_GRAY_DIM)
    add_text_box(slide, "⑥ 保存输出",
                 left=Inches(9.3), top=Inches(2.0), width=Inches(2.5), height=Inches(0.35),
                 font_size=12, color=C_GRAY_DIM)

    # Note
    add_text_box(slide,
                 "35 个测试用例覆盖全部核心逻辑  ·  部分数据源不可用时优雅降级，不崩溃",
                 left=Inches(0.3), top=Inches(6.9), width=Inches(12.7), height=Inches(0.4),
                 font_size=12, color=C_GRAY_DIM)

    slide_number(slide, 8)


def slide_09_visual_rules(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "视觉设计规则 — 专门为金融小白优化", "每个设计决策都有理由")

    rules = [
        (C_GREEN, "图永远比表多",
         "所有数字都有对应的可视化。饼图看 allocation，棒图看盈亏，温度计看情绪。"
         "不会只给你一个数字让你自己猜意思。"),
        (C_RED, "红 / 绿 / 橙全局一致",
         "绿 = 涨 / 好  ·  红 = 跌 / 危险  ·  橙 = 警告。从第一页到最后一页，颜色语义从不变化。"),
        (C_AMBER, "每个数字都有 Tooltip (悬浮解释)",
         "把鼠标悬停在任何数字上，会弹出一句白话解释。VIX 旁边会说: 低于 20 = 市场平静。"),
        (C_BLUE_MID, "底部固定术语面板",
         "每个页面底部都有可折叠的 Glossary，VIX / DXY / F&G / Funding Rate 等全部有定义。"),
        (C_GREEN, "Claude 的输出禁止 jargon",
         "系统 Prompt 强制要求：遇到专业词必须当场括号解释。"
         "比如写'流动性(市场上有多少钱可以买卖)趋紧'而不是只写'流动性趋紧'。"),
    ]

    for i, (color, title, desc) in enumerate(rules):
        y = Inches(2.35 + i * 0.95)
        add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.7), fill=color)
        add_text_box(slide, title,
                     left=Inches(0.6), top=y,
                     width=Inches(12.1), height=Inches(0.38),
                     font_size=18, bold=True, color=C_BG)
        add_text_box(slide, desc,
                     left=Inches(0.6), top=y + Inches(0.38),
                     width=Inches(12.1), height=Inches(0.45),
                     font_size=15, color=C_GRAY_TEXT)

    slide_number(slide, 9)


def slide_10_security(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "安全设计 — 个人工具也应该做对", "四层防护，从 CDN 到本地网络")

    sec_items = [
        (C_GREEN, "SRI — 子资源完整性",
         "HTML 文件里的所有 CDN 脚本（Pico CSS、Chart.js）都带 integrity= 校验码。"
         "如果 CDN 被黑客替换内容，浏览器会拒绝执行，不会运行恶意代码。"),
        (C_BLUE_MID, "DOM API — 防 XSS 注入",
         "Portfolio Explorer 的 WebUI 用 textContent 和 DOM API 渲染用户上传的数据，"
         "从不用 innerHTML。即使 CSV 里包含 <script> 标签，也不会执行。"),
        (C_AMBER, "Loopback Only — 无公网暴露",
         "FastAPI 只绑定 127.0.0.1:8765，不是 0.0.0.0。即使在办公室 WiFi 或公共网络，"
         "其他设备也完全无法访问你的 Portfolio Explorer。"),
        (C_RED, "CSV 永不出本机",
         "portfolios/ 被 .gitignore 排除。git commit / git push 都不会把你的持仓文件带走。"
         "这是代码级别的保证，不只是文档说明。"),
    ]

    for i, (color, title, desc) in enumerate(sec_items):
        y = Inches(2.35 + i * 1.15)
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(1.05),
                 fill=C_WHITE, line_color=color, line_width_pt=2.0)
        add_rect(slide, Inches(0.4), y, Inches(0.18), Inches(1.05), fill=color)
        add_text_box(slide, title,
                     left=Inches(0.7), top=y + Inches(0.08),
                     width=Inches(12.0), height=Inches(0.42),
                     font_size=18, bold=True, color=C_BG)
        add_text_box(slide, desc,
                     left=Inches(0.7), top=y + Inches(0.5),
                     width=Inches(12.0), height=Inches(0.45),
                     font_size=14, color=C_GRAY_TEXT)

    slide_number(slide, 10)


def slide_11_how_to_use(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "怎么用 — 3 种方式，选最顺手的", "从双击启动到全自动 launchd")

    methods = [
        (C_GREEN, "双击启动器（推荐，最简单）",
         "launchers/ 文件夹里有 4 个文件",
         [
             "Daily-Brief.command  /  Daily-Brief.bat",
             "Portfolio-Explorer.command  /  Portfolio-Explorer.bat",
             "Mac 第一次：右键 → 打开 → 打开，之后直接双击",
             "Win 第一次：SmartScreen → 更多信息 → 仍然运行",
         ]),
        (C_BLUE_MID, "命令行（需要打开 Terminal）",
         "python morning_brief.py  /  python webui.py",
         [
             "--no-open      不自动打开浏览器",
             "--commit       把快照 git commit 存档",
             "--commit --push   存档并推送到远程",
             "适合开发者或想要更多控制的场景",
         ]),
        (C_AMBER, "Mac mini launchd 全自动",
         "./scripts/install_launchd.sh",
         [
             "安装后每天 08:00 Europe/Rome 自动跑",
             "起床就能看今天的市场快照",
             "日志存在 .launchd.out.log",
             "卸载：launchctl unload + rm plist",
         ]),
    ]

    for i, (color, title, sub, bullets) in enumerate(methods):
        x = Inches(0.3 + i * 4.35)
        add_rect(slide, x, Inches(2.3), Inches(4.1), Inches(4.8),
                 fill=C_WHITE, line_color=color, line_width_pt=2.5)
        add_rect(slide, x, Inches(2.3), Inches(4.1), Inches(0.6), fill=color)
        add_text_box(slide, title,
                     left=x + Inches(0.1), top=Inches(2.32),
                     width=Inches(3.9), height=Inches(0.52),
                     font_size=15, bold=True, color=C_WHITE)
        add_text_box(slide, sub,
                     left=x + Inches(0.1), top=Inches(3.0),
                     width=Inches(3.9), height=Inches(0.38),
                     font_size=13, color=C_GRAY_DIM,
                     font_name="Courier New")
        for j, b in enumerate(bullets):
            add_text_box(slide, "▸  " + b,
                         left=x + Inches(0.15), top=Inches(3.5 + j * 0.73),
                         width=Inches(3.8), height=Inches(0.6),
                         font_size=14, color=C_GRAY_TEXT)

    slide_number(slide, 11)


def slide_12_roadmap(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide, C_BG_LIGHT)
    header_band(slide, "Roadmap — 还会做什么", "Phase 2 backlog  ·  按优先级排列")

    items = [
        (C_RED, "P2-01", "链上指标（LTH-MVRV / NUPL / SOPR）",
         "打完 Day1Global 剩余 9 个付费指标对应的免费替代，BTC 抄底分从 4 维变 13 维"),
        (C_GREEN, "P2-02", "Fund Finance 行业镜头",
         "追踪电力电子、半导体等 ETF 资金流，为 UniCredit / Polimi 相关岗位申请定制市场视野"),
        (C_BLUE_MID, "P2-03", "PDF 财报上传 → Claude 解读",
         "把一份财报 PDF 拖进来，Claude 提取最重要的 3-5 个要点，白话输出"),
        (C_AMBER, "P2-04", "GitHub Pages 公网快照",
         "--commit --push 自动推到 GH Pages，生成公开链接，发给面试官 / 教授看，零安装"),
        (C_GREEN, "P2-05", "Telegram / Mac 通知",
         "当 BTC 抄底分进入极端区（< 15 或 > 85）时，自动推送到手机"),
    ]

    for i, (color, tag, title, desc) in enumerate(items):
        y = Inches(2.35 + i * 0.95)
        add_rect(slide, Inches(0.4), y, Inches(0.7), Inches(0.7), fill=color)
        add_text_box(slide, tag,
                     left=Inches(0.4), top=y + Inches(0.08),
                     width=Inches(0.7), height=Inches(0.55),
                     font_size=11, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, title,
                     left=Inches(1.25), top=y + Inches(0.05),
                     width=Inches(11.1), height=Inches(0.4),
                     font_size=18, bold=True, color=C_BG)
        add_text_box(slide, desc,
                     left=Inches(1.25), top=y + Inches(0.45),
                     width=Inches(11.1), height=Inches(0.4),
                     font_size=14, color=C_GRAY_TEXT)

    # Footer note
    add_text_box(slide,
                 "完整列表见 README.md  →  Roadmap (P2)  章节",
                 left=Inches(0.4), top=Inches(7.1), width=Inches(12.5), height=Inches(0.3),
                 font_size=12, color=C_GRAY_DIM)

    slide_number(slide, 12)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = new_prs()
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_surface_a(prs)
    slide_05_surface_b(prs)
    slide_06_btc_score(prs)
    slide_07_data_sources(prs)
    slide_08_architecture(prs)
    slide_09_visual_rules(prs)
    slide_10_security(prs)
    slide_11_how_to_use(prs)
    slide_12_roadmap(prs)
    prs.save(str(OUT))
    print(f"Saved {OUT}  ({OUT.stat().st_size // 1024} KB,  {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
